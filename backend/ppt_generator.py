"""
QuantVeil PowerPoint generator — 7 slides.
Uses python-pptx. No template file needed.
"""
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── JP Morgan palette ──────────────────────────────────────────────────────────
NAVY    = RGBColor(0x00, 0x30, 0x87)
NAVY_DK = RGBColor(0x00, 0x1A, 0x4D)
GOLD    = RGBColor(0xB5, 0x88, 0x2C)
GOLD_LT = RGBColor(0xD4, 0xA8, 0x43)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHT  = RGBColor(0xF4, 0xF7, 0xFB)
DTEXT   = RGBColor(0x1A, 0x20, 0x35)
MTEXT   = RGBColor(0x5A, 0x64, 0x78)
LTEXT   = RGBColor(0x84, 0x92, 0xA6)
GREEN   = RGBColor(0x1B, 0x7A, 0x4D)
RED     = RGBColor(0xC0, 0x33, 0x2A)
BGRAY   = RGBColor(0xE8, 0xEC, 0xF2)
BORDER  = RGBColor(0xD4, 0xDC, 0xE8)

W = Inches(13.33)
H = Inches(7.5)
MARGIN = Inches(0.55)
CONTENT_TOP = Inches(1.6)
CONTENT_H   = H - CONTENT_TOP - Inches(0.4)
CONTENT_W   = W - MARGIN * 2


# ── Low-level helpers ──────────────────────────────────────────────────────────

def _prs() -> Presentation:
    p = Presentation()
    p.slide_width  = W
    p.slide_height = H
    return p


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, left, top, width, height, fill_color, border_color=None):
    s = slide.shapes.add_shape(1, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if border_color:
        s.line.color.rgb = border_color
        s.line.width = Pt(1)
    else:
        s.line.color.rgb = fill_color
        s.line.width = Pt(0)
    return s


def _text(slide, left, top, width, height, text, font="Calibri", size=12,
          bold=False, italic=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
    if color is None:
        color = DTEXT
    box = slide.shapes.add_textbox(left, top, width, height)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)[:1000]  # safety truncation
    run.font.name  = font
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def _multiline(slide, left, top, width, height, lines: list[tuple], wrap=True):
    """lines: [(text, font, size, bold, italic, color, align)]"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf  = box.text_frame
    tf.word_wrap = wrap
    first = True
    for (txt, font, size, bold, italic, color, align) in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        run = p.add_run()
        run.text = str(txt)[:500]
        run.font.name  = font
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box


# ── Shared slide components ────────────────────────────────────────────────────

def _header(slide, title: str, subtitle: str = ""):
    # Navy bar
    _rect(slide, 0, 0, W, Inches(1.4), NAVY_DK)
    # Gold accent line
    _rect(slide, 0, Inches(1.4), W, Pt(3), GOLD)
    # Title
    _text(slide, MARGIN, Inches(0.12), W - MARGIN * 2, Inches(0.75),
          title, font="Georgia", size=26, bold=True, color=WHITE)
    # Subtitle
    if subtitle:
        _text(slide, MARGIN, Inches(0.85), W - MARGIN * 2, Inches(0.45),
              subtitle, size=11, color=RGBColor(0xB8, 0xC4, 0xD8))
    # JP Morgan mark in top-right
    _rect(slide, W - Inches(0.45), Inches(0.08), Inches(0.32), Inches(0.32), GOLD)
    _text(slide, W - Inches(0.46), Inches(0.07), Inches(0.34), Inches(0.34),
          "J", font="Georgia", size=13, bold=True, color=NAVY_DK, align=PP_ALIGN.CENTER)


def _footer(slide, domain: str, page: int, total: int = 7):
    _text(slide, MARGIN, H - Inches(0.35), Inches(4), Inches(0.3),
          f"QuantVeil — {domain}", size=8, color=LTEXT)
    _text(slide, W - Inches(1.2), H - Inches(0.35), Inches(0.9), Inches(0.3),
          f"{page} / {total}", size=8, color=LTEXT, align=PP_ALIGN.RIGHT)


def _section_label(slide, left, top, label: str):
    _text(slide, left, top, Inches(5), Inches(0.25),
          label.upper(), size=9, bold=True, color=LTEXT)
    _rect(slide, left, top + Inches(0.25), Inches(2), Pt(2), GOLD)


def _bullet_list(slide, left, top, width, items: list[str],
                 color=DTEXT, size=11, spacing_inches=0.28):
    y = top
    for item in items[:10]:
        # Gold bullet
        _rect(slide, left, y + Inches(0.06), Inches(0.06), Inches(0.06), GOLD)
        _text(slide, left + Inches(0.14), y, width - Inches(0.14), Inches(0.28),
              item[:120], size=size, color=color)
        y += Inches(spacing_inches)
    return y


# ── Slide builders ─────────────────────────────────────────────────────────────

def _slide1(prs, data: dict):
    """Title slide — Company Overview."""
    slide = _blank(prs)
    domain = data.get("domain", "")
    title  = data.get("title", domain)
    desc   = data.get("contacts", {}).get("meta_description", "")[:280] or \
             "Market intelligence report generated by QuantVeil."

    # Full-bleed navy background
    _rect(slide, 0, 0, W, H, NAVY_DK)
    # Gold accent left bar
    _rect(slide, 0, 0, Inches(0.18), H, GOLD)
    # Decorative gold block
    _rect(slide, Inches(0.18), H * 0.45, W - Inches(0.18), Pt(3), GOLD)

    # Company name
    _text(slide, Inches(0.5), Inches(1.5), W - Inches(1), Inches(1.2),
          title[:60], font="Georgia", size=44, bold=True, color=WHITE)
    # Domain badge
    _rect(slide, Inches(0.5), Inches(2.8), Inches(3.2), Inches(0.4),
          RGBColor(0x1B, 0x4F, 0xA8))
    _text(slide, Inches(0.5), Inches(2.78), Inches(3.2), Inches(0.44),
          f"  {domain}", size=12, bold=False, color=RGBColor(0xB8, 0xC4, 0xD8),
          align=PP_ALIGN.LEFT)
    # Description
    _text(slide, Inches(0.5), Inches(3.45), W * 0.6, Inches(1.4),
          desc, size=13, color=RGBColor(0xCC, 0xD4, 0xE0))
    # Report label
    _text(slide, Inches(0.5), H - Inches(1.4), Inches(6), Inches(0.3),
          "QUANTVEIL INTELLIGENCE REPORT", size=9, bold=True,
          color=GOLD, align=PP_ALIGN.LEFT)
    from datetime import date
    _text(slide, Inches(0.5), H - Inches(1.1), Inches(4), Inches(0.3),
          date.today().strftime("%B %d, %Y"), size=9, color=LTEXT)
    # JP Morgan wordmark
    _text(slide, W - Inches(2.2), H - Inches(0.55), Inches(1.8), Inches(0.4),
          "QUANTVEIL", size=9, bold=True,
          color=RGBColor(0x60, 0x70, 0x90), align=PP_ALIGN.RIGHT)


def _slide2(prs, data: dict):
    """Contacts — emails, phones, socials."""
    slide = _blank(prs)
    domain = data.get("domain", "")
    _header(slide, "Contact Intelligence", domain)
    _footer(slide, domain, 2)
    _rect(slide, 0, CONTENT_TOP - Inches(0.05), W, CONTENT_H + Inches(0.05), OFFWHT)

    contacts = data.get("contacts", {})
    emails   = contacts.get("emails",  [])
    phones   = contacts.get("phones",  [])
    socials  = contacts.get("socials", {})

    col_w = (CONTENT_W - Inches(0.3)) / 2
    # Emails column
    _section_label(slide, MARGIN, CONTENT_TOP + Inches(0.1), "Email Addresses")
    y = CONTENT_TOP + Inches(0.5)
    if emails:
        for email in emails[:8]:
            _rect(slide, MARGIN, y, col_w - Inches(0.1), Inches(0.32), WHITE, BORDER)
            _text(slide, MARGIN + Inches(0.1), y + Inches(0.04), col_w - Inches(0.3), Inches(0.24),
                  email, size=11, color=NAVY)
            y += Inches(0.4)
    else:
        _text(slide, MARGIN, y, col_w, Inches(0.3), "No emails detected.", size=11, color=LTEXT)

    # Phones column
    col2_left = MARGIN + col_w + Inches(0.3)
    _section_label(slide, col2_left, CONTENT_TOP + Inches(0.1), "Phone Numbers")
    y2 = CONTENT_TOP + Inches(0.5)
    if phones:
        for phone in phones[:6]:
            _rect(slide, col2_left, y2, col_w, Inches(0.32), WHITE, BORDER)
            _text(slide, col2_left + Inches(0.1), y2 + Inches(0.04), col_w - Inches(0.2), Inches(0.24),
                  phone, size=11, color=NAVY)
            y2 += Inches(0.4)
    else:
        _text(slide, col2_left, y2, col_w, Inches(0.3), "No phones detected.", size=11, color=LTEXT)

    # Socials row at bottom
    if socials:
        sy = max(y, y2) + Inches(0.15)
        _section_label(slide, MARGIN, sy, "Social Media Profiles")
        sx = MARGIN
        for platform, link in list(socials.items())[:8]:
            _rect(slide, sx, sy + Inches(0.35), Inches(1.6), Inches(0.36), NAVY)
            _text(slide, sx + Inches(0.08), sy + Inches(0.37), Inches(1.44), Inches(0.32),
                  f"  {platform.title()}", size=10, bold=True, color=WHITE)
            sx += Inches(1.7)
            if sx > W - Inches(2):
                break


def _slide3(prs, data: dict):
    """Market Position + Reddit Sentiment."""
    slide   = _blank(prs)
    domain  = data.get("domain", "")
    reddit  = data.get("reddit", {})
    wayback = data.get("wayback", {})
    _header(slide, "Market Position & Community Sentiment", domain)
    _footer(slide, domain, 3)
    _rect(slide, 0, CONTENT_TOP - Inches(0.05), W, CONTENT_H + Inches(0.05), OFFWHT)

    summary = reddit.get("ai_summary", "No Reddit data found.")[:500]
    posts   = reddit.get("posts", [])[:5]

    # Left: Reddit summary
    _section_label(slide, MARGIN, CONTENT_TOP + Inches(0.08), "Reddit Community Intelligence")
    _rect(slide, MARGIN, CONTENT_TOP + Inches(0.45),
          Inches(6.5), CONTENT_H - Inches(0.55), WHITE, BORDER)
    _text(slide, MARGIN + Inches(0.12), CONTENT_TOP + Inches(0.55),
          Inches(6.3), CONTENT_H - Inches(0.75),
          summary, size=10, color=DTEXT, wrap=True)

    # Right: Growth + top posts
    rx = MARGIN + Inches(6.8)
    rw = CONTENT_W - Inches(6.8)
    _section_label(slide, rx, CONTENT_TOP + Inches(0.08), "Growth Signal")
    trend = wayback.get("trend", "unknown")
    trend_icon = {"growing": "↑ Growing", "shrinking": "↓ Shrinking", "stable": "→ Stable"}.get(trend, "– Unknown")
    trend_color = {"growing": GREEN, "shrinking": RED, "stable": NAVY}.get(trend, MTEXT)
    _rect(slide, rx, CONTENT_TOP + Inches(0.45), rw, Inches(0.5), WHITE, BORDER)
    _text(slide, rx + Inches(0.1), CONTENT_TOP + Inches(0.5), rw - Inches(0.2), Inches(0.4),
          trend_icon, size=14, bold=True, color=trend_color)

    _section_label(slide, rx, CONTENT_TOP + Inches(1.15), "Top Posts")
    y = CONTENT_TOP + Inches(1.52)
    for post in posts:
        score = post.get("score", 0)
        _rect(slide, rx, y, Inches(0.55), Inches(0.3), GOLD)
        _text(slide, rx, y + Inches(0.04), Inches(0.55), Inches(0.22),
              str(score), size=9, bold=True, color=NAVY_DK, align=PP_ALIGN.CENTER)
        _text(slide, rx + Inches(0.62), y, rw - Inches(0.68), Inches(0.32),
              post.get("title", "")[:55], size=9, color=DTEXT, wrap=False)
        y += Inches(0.38)
        if y > H - Inches(0.6):
            break


def _slide4(prs, data: dict):
    """Pain Points & Opportunities (from market research)."""
    slide  = _blank(prs)
    domain = data.get("domain", "")
    mkt    = data.get("market_research", "")
    _header(slide, "Pain Points & Opportunity Map", domain)
    _footer(slide, domain, 4)
    _rect(slide, 0, CONTENT_TOP - Inches(0.05), W, CONTENT_H + Inches(0.05), OFFWHT)

    # Extract SWOT bullets from market research text
    def _extract_bullets(text, marker, max_items=5):
        """Pull bullets between a marker heading and the next '##' heading."""
        pattern = rf"{re.escape(marker)}.*?\n(.*?)(?=\n##|\Z)"
        m = re.search(pattern, text, re.S | re.I)
        if not m:
            return []
        block = m.group(1)
        bullets = re.findall(r"[W|O|S|T]\d+\.\s+(.+?)(?=\n[WOST]\d|\n##|\Z)", block, re.S)
        if not bullets:
            bullets = [l.strip("- •*").strip() for l in block.splitlines()
                       if l.strip().startswith(("-", "•", "*", "W", "O")) and len(l.strip()) > 10]
        return [b.strip()[:120] for b in bullets[:max_items]]

    weaknesses   = _extract_bullets(mkt, "WEAKNESSES")   or ["Limited data available"]
    opportunities = _extract_bullets(mkt, "OPPORTUNITIES") or ["Limited data available"]

    half = (CONTENT_W - Inches(0.3)) / 2
    # Pain points (left, red accent)
    _rect(slide, MARGIN, CONTENT_TOP + Inches(0.1), Inches(0.06), CONTENT_H - Inches(0.2),
          RED)
    _text(slide, MARGIN + Inches(0.18), CONTENT_TOP + Inches(0.1), half - Inches(0.18), Inches(0.3),
          "PAIN POINTS / WEAKNESSES", size=9, bold=True, color=RED)
    _bullet_list(slide, MARGIN + Inches(0.18), CONTENT_TOP + Inches(0.5),
                 half - Inches(0.25), weaknesses, color=DTEXT, size=11)

    # Opportunities (right, green accent)
    rx = MARGIN + half + Inches(0.3)
    _rect(slide, rx, CONTENT_TOP + Inches(0.1), Inches(0.06), CONTENT_H - Inches(0.2), GREEN)
    _text(slide, rx + Inches(0.18), CONTENT_TOP + Inches(0.1), half - Inches(0.18), Inches(0.3),
          "OPPORTUNITIES", size=9, bold=True, color=GREEN)
    _bullet_list(slide, rx + Inches(0.18), CONTENT_TOP + Inches(0.5),
                 half - Inches(0.25), opportunities, color=DTEXT, size=11)


def _slide5(prs, data: dict):
    """Competitor Landscape."""
    slide  = _blank(prs)
    domain = data.get("domain", "")
    mkt    = data.get("market_research", "")
    _header(slide, "Competitive Landscape", domain)
    _footer(slide, domain, 5)
    _rect(slide, 0, CONTENT_TOP - Inches(0.05), W, CONTENT_H + Inches(0.05), OFFWHT)

    # Extract competitor lines from market research
    comp_section = re.search(r"COMPETITIVE LANDSCAPE(.*?)(?=\n##|\Z)", mkt, re.S | re.I)
    comp_lines = []
    if comp_section:
        for line in comp_section.group(1).splitlines():
            line = line.strip("- •*").strip()
            if len(line) > 10 and not line.startswith("#"):
                comp_lines.append(line[:130])
    if not comp_lines:
        comp_lines = ["No competitor data found in research — check Reddit posts above."]

    _section_label(slide, MARGIN, CONTENT_TOP + Inches(0.1), "Competitors & Alternatives Mentioned")
    y = CONTENT_TOP + Inches(0.52)
    for i, comp in enumerate(comp_lines[:8]):
        row_bg = WHITE if i % 2 == 0 else OFFWHT
        _rect(slide, MARGIN, y, CONTENT_W, Inches(0.4), row_bg, BORDER)
        _rect(slide, MARGIN, y, Inches(0.04), Inches(0.4), NAVY)
        _text(slide, MARGIN + Inches(0.14), y + Inches(0.06), CONTENT_W - Inches(0.2), Inches(0.28),
              comp, size=11, color=DTEXT, wrap=False)
        y += Inches(0.42)


def _slide6(prs, data: dict):
    """News Timeline."""
    slide  = _blank(prs)
    domain = data.get("domain", "")
    news_d = data.get("news", {})
    _header(slide, "News & Media Timeline", domain)
    _footer(slide, domain, 6)
    _rect(slide, 0, CONTENT_TOP - Inches(0.05), W, CONTENT_H + Inches(0.05), OFFWHT)

    news    = news_d.get("news", [])[:8]
    hn_posts = news_d.get("hn_posts", [])[:5]

    # Left: News
    half = (CONTENT_W - Inches(0.3)) / 2
    _section_label(slide, MARGIN, CONTENT_TOP + Inches(0.08), "Google News")
    y = CONTENT_TOP + Inches(0.48)
    for item in news[:6]:
        date  = item.get("date", "")[:12]
        src   = item.get("source", "")[:20]
        title = item.get("title", "")[:70]
        _text(slide, MARGIN, y, half, Inches(0.18),
              f"{date}  ·  {src}", size=8, color=LTEXT)
        _text(slide, MARGIN, y + Inches(0.17), half, Inches(0.25),
              title, size=10, color=DTEXT)
        _rect(slide, MARGIN, y + Inches(0.43), half, Pt(1), BORDER)
        y += Inches(0.5)

    # Right: HN
    rx = MARGIN + half + Inches(0.3)
    _section_label(slide, rx, CONTENT_TOP + Inches(0.08), "Hacker News")
    y2 = CONTENT_TOP + Inches(0.48)
    for post in hn_posts:
        pts   = post.get("points", 0)
        title = post.get("title", "")[:65]
        _rect(slide, rx, y2 + Inches(0.04), Inches(0.52), Inches(0.28), GOLD)
        _text(slide, rx, y2 + Inches(0.05), Inches(0.52), Inches(0.26),
              f"{pts}pts", size=8, bold=True, color=NAVY_DK, align=PP_ALIGN.CENTER)
        _text(slide, rx + Inches(0.6), y2, half - Inches(0.65), Inches(0.36),
              title, size=10, color=DTEXT, wrap=False)
        y2 += Inches(0.48)


def _slide7(prs, data: dict):
    """Cold Outreach Strategy."""
    slide  = _blank(prs)
    domain = data.get("domain", "")
    cold   = data.get("cold_outreach", "")
    _header(slide, "Cold Outreach Strategy", domain)
    _footer(slide, domain, 7)
    _rect(slide, 0, CONTENT_TOP - Inches(0.05), W, CONTENT_H + Inches(0.05), OFFWHT)

    if not cold:
        _text(slide, MARGIN, CONTENT_TOP + Inches(0.5), CONTENT_W, Inches(0.4),
              "Enable AI analysis to generate outreach strategy.", size=12, color=LTEXT)
        return

    # Opening line — highlighted box
    opening_m = re.search(r"(?:Opening Line|Cold Outreach)[:\s]*\"?(.+?)\"?\n", cold, re.I)
    opening   = opening_m.group(1).strip()[:180] if opening_m else cold[:180]

    _text(slide, MARGIN, CONTENT_TOP + Inches(0.12), CONTENT_W, Inches(0.25),
          "OPENING LINE", size=9, bold=True, color=LTEXT)
    _rect(slide, MARGIN, CONTENT_TOP + Inches(0.38), CONTENT_W, Inches(0.75), WHITE, GOLD)
    _rect(slide, MARGIN, CONTENT_TOP + Inches(0.38), Inches(0.06), Inches(0.75), GOLD)
    _text(slide, MARGIN + Inches(0.16), CONTENT_TOP + Inches(0.48),
          CONTENT_W - Inches(0.24), Inches(0.55),
          f'"{opening}"', font="Georgia", size=12, italic=True, color=NAVY)

    # Strategy points
    _text(slide, MARGIN, CONTENT_TOP + Inches(1.28), CONTENT_W, Inches(0.25),
          "STRATEGY & TALKING POINTS", size=9, bold=True, color=LTEXT)

    # Extract bullets from cold outreach text
    bullets = []
    for line in cold.splitlines():
        ln = line.strip("- •*#").strip()
        if len(ln) > 15 and not ln.startswith("**") and "Opening" not in ln and "Cold" not in ln:
            bullets.append(ln[:130])
    if not bullets:
        bullets = [cold[i:i+130] for i in range(0, min(len(cold), 650), 130)]

    _bullet_list(slide, MARGIN, CONTENT_TOP + Inches(1.6), CONTENT_W,
                 bullets[:7], color=DTEXT, size=11, spacing_inches=0.42)


# ── Public API ─────────────────────────────────────────────────────────────────

def generate_ppt(data: dict, output_dir: str = ".") -> str:
    prs = _prs()
    _slide1(prs, data)
    _slide2(prs, data)
    _slide3(prs, data)
    _slide4(prs, data)
    _slide5(prs, data)
    _slide6(prs, data)
    _slide7(prs, data)

    domain   = re.sub(r"[^a-z0-9]", "_", data.get("domain", "report"))
    filename = f"intel_{domain}.pptx"
    path     = os.path.join(output_dir, filename)
    prs.save(path)
    return path
